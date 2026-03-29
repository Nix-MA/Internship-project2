"""
ingestion/extractor.py — Document text extraction for all 18 supported file types.

Output contract for extract_document():
    {
        "type": "text" | "image" | "audio" | "video" | "archive",
        "content": <str>,
        "metadata": {
            "filename": <str>,
            "size_kb": <float>,
            "extension": <str>,
            "num_pages": <int>   # PDF only
        }
    }
"""

import os
import csv
import json as jsonlib
import zipfile
import base64
import requests as _requests
from src.utils.logger import get_logger

try:
    from PIL import Image
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False

VISION_MODELS = ["llama3.2-vision", "moondream"]
VISION_TIMEOUT = 120  # seconds

MAX_FILE_SIZE_MB = 10
MAX_UNCOMPRESSED_RATIO = 100 # For zip bomb detection

# ── File type detection ────────────────────────────────────────────────────────

EXT_TYPE_MAP = {
    ".pdf":  "pdf",
    ".txt":  "txt",
    ".md":   "txt",
    ".docx": "docx",
    ".csv":  "csv",
    ".json": "json",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".png":  "image",
    ".jpg":  "image",
    ".jpeg": "image",
    ".webp": "image",
    ".gif":  "image",
    ".bmp":  "image",
    ".tiff": "image",
    ".mp3":  "audio",
    ".wav":  "audio",
    ".mp4":  "video",
    ".mov":  "video",
    ".avi":  "video",
    ".mkv":  "video",
    ".zip":  "archive",
}


def detect_file_type(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    return EXT_TYPE_MAP.get(ext, "unknown")


# ── Per-type extractors ────────────────────────────────────────────────────────

def _extract_pdf(filepath: str) -> tuple[str, dict]:
    import pdfplumber
    text = ""
    num_pages = 0
    with pdfplumber.open(filepath) as pdf:
        num_pages = len(pdf.pages)
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text.strip(), {"num_pages": num_pages}


def _extract_txt(filepath: str) -> tuple[str, dict]:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip(), {}


def _extract_docx(filepath: str) -> tuple[str, dict]:
    try:
        from docx import Document
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs), {"num_paragraphs": len(paragraphs)}
    except Exception as e:
        return f"[DOCX extraction failed: {e}]", {}


def _extract_csv(filepath: str) -> tuple[str, dict]:
    try:
        lines = []
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                lines.append(", ".join(str(c) for c in row))
        return "\n".join(lines[:300]), {"num_rows": len(lines)}
    except Exception as e:
        return f"[CSV extraction failed: {e}]", {}


def _extract_json(filepath: str) -> tuple[str, dict]:
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            data = jsonlib.load(f)
        text = jsonlib.dumps(data, indent=2)
        return text[:8000], {"num_keys": len(data) if isinstance(data, dict) else len(data)}
    except Exception as e:
        return f"[JSON extraction failed: {e}]", {}


def _extract_xlsx(filepath: str) -> tuple[str, dict]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                row_str = ", ".join(str(c) for c in row if c is not None)
                if row_str:
                    lines.append(row_str)
        return "\n".join(lines[:400]), {"sheets": wb.sheetnames}
    except Exception as e:
        return f"[XLSX extraction failed: {e}]", {}


def _extract_pptx(filepath: str) -> tuple[str, dict]:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"[Slide {i + 1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines), {"num_slides": len(prs.slides)}
    except Exception as e:
        return f"[PPTX extraction failed: {e}]", {}


def _extract_zip(filepath: str) -> tuple[str, dict]:
    filename = os.path.basename(filepath)
    size_kb = round(os.path.getsize(filepath) / 1024, 1)
    
    try:
        with zipfile.ZipFile(filepath, 'r') as zf:
            total_uncompressed = 0
            for info in zf.infolist():
                if info.filename.endswith('.zip'):
                    get_logger('ingestion').warning(f"Nested ZIP detected in {filename}")
                    return "[Archive rejected: nested compression detected]", {"size_kb": size_kb}
                total_uncompressed += info.file_size
                
                # Check ratio for files > 100KB
                if info.compress_size > 100 * 1024 and info.compress_size > 0:
                    ratio = info.file_size / info.compress_size
                    if ratio > MAX_UNCOMPRESSED_RATIO:
                        get_logger('ingestion').warning(f"Zip bomb detected in {filename} (ratio {ratio})")
                        return "[Archive rejected: potentially malicious compression ratio]", {"size_kb": size_kb}
                        
            return f"[Archive {filename} validated. Contains {len(zf.infolist())} files. Uncompressed size: {total_uncompressed / 1024:.1f} KB]", {"size_kb": size_kb}
    except zipfile.BadZipFile:
        get_logger('ingestion').error(f"Bad zip file: {filename}")
        return f"[Archive rejected: corrupted or invalid format]", {"size_kb": size_kb}
    except Exception as e:
        get_logger('ingestion').error(f"Zip inspection failed: {filename}", exc_info=True)
        return f"[Archive inspection failed]", {"size_kb": size_kb}

def _describe_image_with_vision(filepath: str, model: str) -> str | None:
    """
    Send an image to a vision model via Ollama and get a detailed textual
    description of its content.  Returns None if the model is unavailable.
    """
    try:
        with open(filepath, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")

        payload = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": (
                        "Transcribe all text in this image precisely word-for-word. "
                        "If the image contains data, tables, or facts, state them directly. "
                        "CRITICAL RULES: \n"
                        "1. DO NOT mention what the image is (e.g., skip phrases like 'This is a worksheet', 'The image shows', 'In this picture').\n"
                        "2. NEVER use the words 'image', 'picture', 'worksheet', 'page', 'document', 'screenshot', or 'file'.\n"
                        "3. Just output the raw information and text contained within it."
                    ),
                    "images": [b64],
                }
            ],
            "stream": False,
        }
        resp = _requests.post(
            "http://127.0.0.1:11434/api/chat",
            json=payload,
            timeout=VISION_TIMEOUT,
        )
        resp.raise_for_status()
        description = resp.json().get("message", {}).get("content", "").strip()
        if description:
            get_logger('ingestion').info(f"Vision model '{model}' described image ({len(description)} chars)")
            return description
    except Exception as e:
        get_logger('ingestion').warning(f"Vision model '{model}' failed for image: {e}")
    return None


def _extract_media(filepath: str, media_type: str) -> tuple[str, dict]:
    """
    For images: send to llama3.2-vision for real content extraction.
    For audio/video: return empty string (no transcript = no content).
    For archive: return empty string (no extractable content).

    CRITICAL: Never return format/type/size text — that causes questions
    about file formats rather than real content.
    """
    filename = os.path.basename(filepath)
    stem = os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title()
    size_kb = round(os.path.getsize(filepath) / 1024, 1)

    if media_type == "image":
        ocr_text = ""
        # 1. Always try OCR first. If the image is a worksheet, document screenshot, or
        # infographic, OCR will extract the raw text (which is what we actually want to test on).
        if PYTESSERACT_AVAILABLE:
            try:
                img = Image.open(filepath)
                # Convert to RGB to avoid issues with some formats (like RGBA pngs or certain GIFs)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                ocr_text = pytesseract.image_to_string(img).strip()
            except Exception as e:
                get_logger('ingestion').warning(f"OCR failed for {filename}: {e}")

        # If OCR extracted substantial text (> 30 words), treat this as a text document screenshot.
        # Skip the vision model entirely to prevent meta-descriptions ("This is a worksheet...")
        # which cause the LLM to generate questions about the FORMAT rather than the CONTENT.
        if ocr_text and len(ocr_text.split()) > 30:
            get_logger('ingestion').info(f"OCR found {len(ocr_text.split())} words. Skipping vision model.")
            return ocr_text, {"size_kb": size_kb, "ocr": True, "vision": False}

        # 2. If OCR didn't find much text (it's a photo, diagram, or chart), 
        # fall back to the vision models to describe the visual elements.
        for model in VISION_MODELS:
            vision_desc = _describe_image_with_vision(filepath, model)
            if vision_desc:
                # Combine any sparse text (labels) with the visual description
                combined = vision_desc
                if ocr_text:
                    combined += f"\n\nText found in image:\n{ocr_text}"
                return combined, {"size_kb": size_kb, "vision": True, "vision_model": model, "ocr": bool(ocr_text)}

        # Neither model nor OCR worked — return empty
        get_logger('ingestion').warning(f"No vision model or OCR text available for {filename}.")
        return "", {"size_kb": size_kb, "vision": False, "ocr": False, "status": "insufficient_content"}

    # Audio/video/archive: without a transcript there is no extractable semantic content.
    # Return empty string — the caller will detect this as insufficient_content.
    get_logger('ingestion').info(
        f"{media_type.title()} file '{filename}' has no transcript/text. "
        "Returning empty content; no questions will be generated from it."
    )
    return "", {"size_kb": size_kb, "status": "insufficient_content"}


# ── Master router ──────────────────────────────────────────────────────────────

def extract_document(filepath: str) -> dict:
    """
    Extract content from any supported file type.

    Returns:
        {
            "type": str,
            "content": str,
            "metadata": dict
        }
    """
    file_type = detect_file_type(filepath)
    filename = os.path.basename(filepath)
    size_kb = round(os.path.getsize(filepath) / 1024, 1)
    ext = os.path.splitext(filename)[1].lower()

    base_meta = {"filename": filename, "size_kb": size_kb, "extension": ext}
    
    if size_kb > MAX_FILE_SIZE_MB * 1024:
        get_logger('ingestion').warning(f"File {filename} ({size_kb}KB) exceeded max size limit of {MAX_FILE_SIZE_MB}MB.")
        return {
            "type": "text", 
            "content": f"[File skipped: exceeds {MAX_FILE_SIZE_MB}MB strict limit]", 
            "metadata": base_meta
        }

    try:
        if file_type == "pdf":
            content, extra = _extract_pdf(filepath)
            doc_category = "text"
        elif file_type == "txt":
            content, extra = _extract_txt(filepath)
            doc_category = "text"
        elif file_type == "docx":
            content, extra = _extract_docx(filepath)
            doc_category = "text"
        elif file_type == "csv":
            content, extra = _extract_csv(filepath)
            doc_category = "text"
        elif file_type == "json":
            content, extra = _extract_json(filepath)
            doc_category = "text"
        elif file_type == "xlsx":
            content, extra = _extract_xlsx(filepath)
            doc_category = "text"
        elif file_type == "pptx":
            content, extra = _extract_pptx(filepath)
            doc_category = "text"
        elif file_type == "image":
            content, extra = _extract_media(filepath, "image")
            doc_category = "image"
        elif file_type == "audio":
            content, extra = _extract_media(filepath, "audio")
            doc_category = "audio"
        elif file_type == "video":
            content, extra = _extract_media(filepath, "video")
            doc_category = "video"
        elif file_type == "archive":
            # Archives have no extractable semantic content without decompression+parsing.
            # Return empty — downstream will skip gracefully.
            get_logger('ingestion').info(f"Archive '{filename}' skipped: no semantic content extractable.")
            content = ""
            extra = {"status": "insufficient_content"}
            doc_category = "archive"
        else:
            # Last resort: try plain text, but check for binary first
            try:
                with open(filepath, 'rb') as bf:
                    chunk = bf.read(1024)
                    if b'\0' in chunk:
                        # Binary file with no text — return empty, not a format description
                        get_logger('ingestion').warning(f"Unsupported binary file: {filename}. Returning empty content.")
                        content = ""
                        extra = {"status": "insufficient_content"}
                        doc_category = "text"
                    else:
                        content, extra = _extract_txt(filepath)
                        doc_category = "text"
            except Exception:
                get_logger('ingestion').warning(f"Could not read unknown file: {filename}. Returning empty content.")
                content = ""
                extra = {"status": "insufficient_content"}
                doc_category = "text"

    except Exception as e:
        get_logger('ingestion').error(f"Fatal extraction error for {filename}", exc_info=True)
        # Return empty string — never return error text that could seed hallucinated questions
        content = ""
        extra = {"status": "error", "error": str(e)}
        doc_category = "text"

    return {
        "type": doc_category,
        "content": content,
        "metadata": {**base_meta, **extra},
    }
